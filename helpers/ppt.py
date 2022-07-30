from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.dml import MSO_LINE

def create_slide_linechart(slide, sub_df, gold_code, page_num):
    categories = sub_df['date'].tolist()
    buying_prices = tuple(sub_df['avg_buying_price'].tolist())
    selling_prices = tuple(sub_df['avg_selling_price'].tolist())
    start_date = categories[0]
    end_date = categories[-1]
    
    # define chart data 
    chart_data = ChartData()
    chart_data.categories = categories 
    chart_data.add_series('Buying Price', buying_prices)
    chart_data.add_series('Selling Price', selling_prices)
    
    # add chart to slide 
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(11), Inches(4.5)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    plot = chart.plots[0]
    
    # add heading
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        text_frame.text = f'CODE = {gold_code.upper()}'
        break

    # add title
    chart_title = slide.shapes.add_textbox(Inches(3.0), Inches(0.8), Inches(7), Inches(0.4))
    tf = chart_title.text_frame
    p = tf.add_paragraph()
    p.text = "Gold Price of "
    first_run = p.add_run()
    first_run.text = f"{gold_code}" 
    first_run.font.bold = True
    first_run.font.color.rgb = RGBColor(0, 0, 255)
    second_run = p.add_run()
    second_run.text = f" in last 14 Days (from {start_date} to {end_date})"
    
    # add page number
    page_number = slide.shapes.add_textbox(Inches(0.25), Inches(7), Inches(0.7), Inches(0.7))
    page_number.text = f"Page {str(page_num)}"
    page_number.text_frame.paragraphs[0].font.size = Pt(9)

    # Customize series
    plot.has_data_labels = True

    for series in plot.series:
        if series.name == 'Buying Price':
            series_color = RGBColor(255, 0, 0)
        elif series.name == 'Selling Price':
            series_color = RGBColor(0, 255, 0) 
        for point, value in zip(series.points, series.values):
            data_label = point.data_label
            data_label.font.color.rgb = series_color
            data_label.font.size = Pt(9.5)
            data_label.position = XL_LABEL_POSITION.ABOVE
            
        line = series.format.line
        line.color.rgb = series_color
        line.dash_style = MSO_LINE.DASH
        line.width = Pt(1.5)
        series.smooth = True

    # Customize series
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = False
    value_axis.axis_title.text_frame.text = 'Price'
    value_axis.tick_labels.font.size= Pt(10)

    category_axis = chart.category_axis
    category_axis.has_minor_gridlines = True
    category_axis.axis_title.text_frame.text = 'Date'
    category_axis.tick_labels.font.size= Pt(10)

    return slide