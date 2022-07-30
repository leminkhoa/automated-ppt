import pendulum
import pandas as pd
import numpy as np
from pptx import Presentation
from helpers.ppt import *

# Load data
data = pd.read_csv('data/staging_data.csv')
gold_codes = data['gold_code'].unique().astype('str').tolist()
data = data.replace(np.nan, None)

# create presentation with 1 slide ------
prs = Presentation('ppt-template/template.pptx')
today = pendulum.now("Asia/Bangkok").to_date_string()

# Intro
intro = prs.slides[0]
for shape in intro.shapes:
    if not shape.has_text_frame:
        continue
    text_frame = shape.text_frame
    if text_frame.text == "Date: #":
        text_frame.text = text_frame.text.replace('#', today)

page_num = 1
for gold_code in gold_codes:
    sub_df = data.query(f"gold_code=='{gold_code}'")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    create_slide_linechart(slide, sub_df, gold_code, page_num)
    page_num += 1

prs.save('output/report.pptx')